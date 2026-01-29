"""
PDF/PPTX to PPTX Converter with AI-Generated Speaker Notes
Converts PDF pages or existing PPTX slides to PowerPoint with AI-generated presenter notes.
"""

import os
import sys
import io
import uuid
import base64
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF
from PIL import Image
from google import genai
from google.genai import types
from dotenv import load_dotenv

# Load environment variables
load_dotenv()


def render_pdf_page_as_image(pdf_path, page_num, dpi=300):
    """
    Render a PDF page as a high-resolution image.
    
    Args:
        pdf_path (str): Path to PDF file
        page_num (int): Page number (0-indexed)
        dpi (int): Resolution
    
    Returns:
        PIL Image
    """
    doc = fitz.open(pdf_path)
    page = doc[page_num]
    
    # Render page to image
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat)
    
    # Convert to PIL Image
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    
    doc.close()
    return img


def generate_speaker_notes(image, api_key, note_style="standard", note_tone="professional"):
    """
    Use Google Gemini to generate speaker notes for a slide.
    
    Args:
        image: PIL Image of the slide
        api_key: Google AI API key
        note_style: Style of notes - 'brief', 'standard', or 'detailed'
        note_tone: Tone of notes - 'professional', 'casual', 'academic', or 'persuasive'
    
    Returns:
        str: Generated speaker notes
    """
    # Convert image to bytes
    img_byte_arr = io.BytesIO()
    image.save(img_byte_arr, format='PNG')
    img_byte_arr = img_byte_arr.getvalue()
    
    # Initialize Gemini client
    client = genai.Client(api_key=api_key)
    
    # Configure style based on selection
    style_configs = {
        'brief': {
            'duration': '20-30 seconds',
            'detail': 'concise and to the point',
            'sentences': '2-4 sentences'
        },
        'standard': {
            'duration': '45-60 seconds',
            'detail': 'clear and comprehensive',
            'sentences': '4-6 sentences'
        },
        'detailed': {
            'duration': '90-120 seconds',
            'detail': 'thorough and detailed with context and examples',
            'sentences': '8-12 sentences'
        }
    }
    
    # Configure tone based on selection
    tone_configs = {
        'professional': 'professional and business-appropriate',
        'casual': 'friendly and conversational',
        'academic': 'scholarly and research-oriented',
        'persuasive': 'compelling and convincing',
        'enthusiastic': 'energetic and passionate',
        'storytelling': 'narrative-driven and engaging with stories',
        'technical': 'precise and technically detailed',
        'inspirational': 'motivational and uplifting',
        'educational': 'clear and instructive for learning',
        'enthusiastic': 'energetic and passionate',
        'storytelling': 'narrative-driven and engaging with stories',
        'technical': 'precise and technically detailed',
        'inspirational': 'motivational and uplifting',
        'educational': 'clear and instructive for learning'
    }
    
    style_config = style_configs.get(note_style, style_configs['standard'])
    tone_description = tone_configs.get(note_tone, tone_configs['professional'])
    
    prompt = f"""Analyze this presentation slide and write exactly what the presenter should say when presenting this slide.

Write a natural, conversational script that:
- Takes approximately {style_config['duration']} to speak
- Is {style_config['detail']}
- Contains {style_config['sentences']}
- Uses a {tone_description} tone
- Flows smoothly and sounds natural when spoken aloud
- Explains what's on the slide clearly
- Is written in first person (as if you are the presenter)
- Uses simple, clear language
- Includes no markdown formatting, bullets, or special characters
- Is just plain text that can be read directly

Write ONLY the spoken words - nothing else. No labels, no sections, no formatting.
Just write what needs to be said, as if you're speaking directly to the audience."""

    try:
        model_name = os.environ.get('GEMINI_MODEL', 'gemini-2.0-flash-exp')
        response = client.models.generate_content(
            model=model_name,
            contents=[
                types.Part.from_bytes(
                    data=img_byte_arr,
                    mime_type='image/png'
                ),
                types.Part.from_text(text=prompt)
            ]
        )
        
        notes = response.text.strip()
        return notes
        
    except Exception as e:
        print(f"    Warning: Failed to generate notes: {e}")
        return "Speaker notes could not be generated for this slide."


def render_slide_as_image(prs, slide_idx, temp_dir="temp_slides"):
    """
    Create a visual representation of a PPTX slide with text overlay.
    
    Args:
        prs: Presentation object
        slide_idx (int): Index of slide to render
        temp_dir (str): Directory for temporary files
    
    Returns:
        PIL Image or None
    """
    try:
        from PIL import ImageDraw, ImageFont
        
        source_slide = prs.slides[slide_idx]
        
        # Get slide dimensions
        slide_width_inches = prs.slide_width.inches
        slide_height_inches = prs.slide_height.inches
        
        # Create a canvas with slide dimensions at 150 DPI
        dpi = 150
        img_width = int(slide_width_inches * dpi)
        img_height = int(slide_height_inches * dpi)
        
        # Create white background
        slide_image = Image.new('RGB', (img_width, img_height), color='white')
        draw = ImageDraw.Draw(slide_image)
        
        # Try to load a font (fallback to default if not available)
        try:
            font = ImageFont.truetype("arial.ttf", 20)
            font_small = ImageFont.truetype("arial.ttf", 16)
        except:
            font = ImageFont.load_default()
            font_small = ImageFont.load_default()
        
        has_content = False
        
        # Process all shapes in the slide
        for shape in source_slide.shapes:
            try:
                # Handle pictures
                if shape.shape_type == 13:  # Picture type
                    image = shape.image
                    image_bytes = image.blob
                    shape_image = Image.open(io.BytesIO(image_bytes))
                    
                    # Calculate position and size
                    left = int((shape.left / prs.slide_width) * img_width)
                    top = int((shape.top / prs.slide_height) * img_height)
                    width = int((shape.width / prs.slide_width) * img_width)
                    height = int((shape.height / prs.slide_height) * img_height)
                    
                    # Resize and paste
                    shape_image = shape_image.resize((width, height), Image.LANCZOS)
                    slide_image.paste(shape_image, (left, top))
                    has_content = True
                
                # Handle text boxes and shapes with text
                elif hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Calculate position and size
                    left = int((shape.left / prs.slide_width) * img_width)
                    top = int((shape.top / prs.slide_height) * img_height)
                    width = int((shape.width / prs.slide_width) * img_width)
                    height = int((shape.height / prs.slide_height) * img_height)
                    
                    # Draw a light background for text
                    draw.rectangle([left, top, left + width, top + height], 
                                 fill='#f0f0f0', outline='#cccccc')
                    
                    # Draw text (simplified - just the first 500 chars)
                    text_display = text[:500]
                    # Wrap text to fit width
                    words = text_display.split()
                    lines = []
                    current_line = ""
                    
                    for word in words:
                        test_line = current_line + " " + word if current_line else word
                        # Rough estimate - 10 pixels per char
                        if len(test_line) * 10 < width:
                            current_line = test_line
                        else:
                            if current_line:
                                lines.append(current_line)
                            current_line = word
                    
                    if current_line:
                        lines.append(current_line)
                    
                    # Draw the lines
                    y_offset = top + 10
                    for line in lines[:10]:  # Limit to 10 lines
                        draw.text((left + 10, y_offset), line, fill='black', font=font_small)
                        y_offset += 25
                    
                    has_content = True
                
                # Handle tables
                elif hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    left = int((shape.left / prs.slide_width) * img_width)
                    top = int((shape.top / prs.slide_height) * img_height)
                    width = int((shape.width / prs.slide_width) * img_width)
                    height = int((shape.height / prs.slide_height) * img_height)
                    
                    # Draw table border
                    draw.rectangle([left, top, left + width, top + height], 
                                 outline='#666666', width=2)
                    
                    # Draw simplified table representation
                    row_height = height // len(table.rows) if len(table.rows) > 0 else 30
                    y_pos = top
                    
                    for row in table.rows[:5]:  # Limit to first 5 rows
                        x_pos = left
                        col_width = width // len(row.cells) if len(row.cells) > 0 else 100
                        
                        for cell in row.cells[:5]:  # Limit to first 5 columns
                            cell_text = cell.text.strip()[:30]  # Limit text length
                            if cell_text:
                                draw.text((x_pos + 5, y_pos + 5), cell_text, 
                                        fill='black', font=font_small)
                            x_pos += col_width
                        
                        y_pos += row_height
                    
                    has_content = True
                    
            except Exception as e:
                # Skip shapes that cause errors
                continue
        
        # Return the image if we found any content
        if has_content:
            return slide_image
        else:
            return None
            
    except Exception as e:
        print(f"    Warning: Could not render slide as image: {str(e)[:100]}")
        return None


def add_notes_to_pptx(input_pptx, output_pptx, api_key):
    """
    Add AI-generated speaker notes to an existing PPTX file.
    Extracts text and images directly from slides for analysis.
    
    Args:
        input_pptx (str): Path to input PPTX file
        output_pptx (str): Path for output PPTX file
        api_key (str): Google AI API key
    
    Returns:
        str: Path to created PPTX file
    """
    print(f"Adding speaker notes to: {input_pptx}")
    print(f"Output PPTX: {output_pptx}")
    print(f"Using Google Gemini for speaker notes generation...\n")
    
    # Open existing presentation
    prs = Presentation(input_pptx)
    num_slides = len(prs.slides)
    print(f"Total slides: {num_slides}\n")
    
    # Process each slide
    for idx, slide in enumerate(prs.slides):
        print(f"{'='*60}")
        print(f"Processing slide {idx + 1}/{num_slides}...")
        print(f"{'='*60}")
        
        notes = None
        
        # Extract text content from slide first
        slide_text = []
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                # Also check for text in tables
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                slide_text.append(text)
        except Exception as e:
            print(f"  Warning: Error extracting text: {str(e)[:100]}")
        
        combined_text = "\n".join(slide_text)
        
        # Try to create a visual representation
        print("  [1/2] Creating slide visual representation...")
        slide_image = render_slide_as_image(prs, idx)
        
        # Strategy: Use image if available for better context, otherwise use text
        if slide_image:
            try:
                print(f"  [2/2] Generating speaker notes with AI (visual analysis with {len(combined_text)} chars of text)...")
                notes = generate_speaker_notes(slide_image, api_key)
            except Exception as e:
                print(f"  Warning: Image-based analysis failed: {str(e)[:100]}")
                notes = None
        
        # Fallback to text-only approach if image analysis failed or no image
        if notes is None and combined_text:
            print(f"  [2/2] Generating speaker notes with AI (text-based, {len(combined_text)} chars)...")
            try:
                notes = generate_notes_from_text(combined_text, api_key)
            except Exception as e:
                print(f"  Error generating notes: {str(e)[:100]}")
                notes = "This slide contains content but speaker notes could not be generated. Please review and add custom notes."
        
        # If we still don't have notes, provide a default message
        if notes is None:
            print("  Warning: No content found in slide")
            notes = "This slide appears to be empty or contains only visual elements without text. Please review and add custom speaker notes as needed."
        
        # Add notes to slide
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
            
            print(f"  ✓ Slide {idx + 1} completed with notes")
            print(f"  Notes preview: {notes[:100]}...\n")
        except Exception as e:
            print(f"  Error adding notes to slide: {str(e)[:100]}\n")
    
    # Save presentation
    print(f"{'='*60}")
    print(f"Saving presentation: {output_pptx}")
    prs.save(output_pptx)
    print("✓ Conversion completed successfully!")
    print("✓ Each slide has AI-generated speaker notes!")
    print(f"{'='*60}\n")
    
    return output_pptx


def generate_notes_from_text(slide_text, api_key, note_style="standard", note_tone="professional"):
    """
    Generate speaker notes from slide text content.
    
    Args:
        slide_text (str): Text content from slide
        api_key (str): Google AI API key
        note_style: Style of notes - 'brief', 'standard', or 'detailed'
        note_tone: Tone of notes - 'professional', 'casual', 'academic', or 'persuasive'
    
    Returns:
        str: Generated speaker notes
    """
    if not slide_text or not slide_text.strip():
        return "This slide appears to contain visual content without text. Please review the slide and add appropriate speaker notes."
    
    client = genai.Client(api_key=api_key)
    
    # Configure style based on selection
    style_configs = {
        'brief': {
            'duration': '20-30 seconds',
            'detail': 'concise and to the point',
            'sentences': '2-4 sentences'
        },
        'standard': {
            'duration': '45-60 seconds',
            'detail': 'clear and comprehensive',
            'sentences': '4-6 sentences'
        },
        'detailed': {
            'duration': '90-120 seconds',
            'detail': 'thorough and detailed with context and examples',
            'sentences': '8-12 sentences'
        }
    }
    
    # Configure tone based on selection
    tone_configs = {
        'professional': 'professional and business-appropriate',
        'casual': 'friendly and conversational',
        'academic': 'scholarly and research-oriented',
        'persuasive': 'compelling and convincing',
        'enthusiastic': 'energetic and passionate',
        'storytelling': 'narrative-driven and engaging with stories',
        'technical': 'precise and technically detailed',
        'inspirational': 'motivational and uplifting',
        'educational': 'clear and instructive for learning'
    }
    
    style_config = style_configs.get(note_style, style_configs['standard'])
    tone_description = tone_configs.get(note_tone, tone_configs['professional'])
    
    prompt = f"""Based on this slide content, write exactly what the presenter should say when presenting this slide.

Slide content:
{slide_text}

Write a natural, conversational script that:
- Takes approximately {style_config['duration']} to speak
- Is {style_config['detail']}
- Contains {style_config['sentences']}
- Uses a {tone_description} tone
- Flows smoothly and sounds natural when spoken aloud
- Explains the content clearly
- Is written in first person (as if you are the presenter)
- Uses simple, clear language
- Includes no markdown formatting, bullets, or special characters
- Is just plain text that can be read directly
- Expands on the bullet points or headings with context and explanation

Write ONLY the spoken words - nothing else. No labels, no sections, no formatting.
Just write what needs to be said, as if you're speaking directly to the audience."""

    try:
        model_name = os.environ.get('GEMINI_MODEL', 'gemini-2.0-flash-exp')
        response = client.models.generate_content(
            model=model_name,
            contents=[
                types.Part.from_text(text=prompt)
            ]
        )
        
        notes = response.text.strip()
        return notes
        
    except Exception as e:
        print(f"    Warning: Failed to generate notes: {e}")
        return "Speaker notes could not be generated for this slide."


def pdf_to_pptx_with_notes(pdf_path, output_pptx=None, dpi=200, api_key=None):
    """
    Convert PDF to PPTX with each page as an image slide plus AI-generated speaker notes.
    
    Args:
        pdf_path (str): Path to input PDF file
        output_pptx (str): Path for output PPTX file
        dpi (int): Resolution for PDF rendering
        api_key (str): Google AI API key
    
    Returns:
        str: Path to created PPTX file
    """
    
    # Validate input
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
    
    # Get API key
    if api_key is None:
        api_key = os.environ.get('GOOGLE_API_KEY')
    
    if not api_key:
        raise ValueError("Google API key required. Set GOOGLE_API_KEY in .env file.")
    
    # Set output filename
    if output_pptx is None:
        output_pptx = str(Path(pdf_path).stem + '_with_notes.pptx')
    
    print(f"Converting PDF: {pdf_path}")
    print(f"Output PPTX: {output_pptx}")
    print(f"DPI: {dpi}")
    print(f"Using Google Gemini for speaker notes generation...\n")
    
    # Get page count
    doc = fitz.open(pdf_path)
    num_pages = len(doc)
    doc.close()
    print(f"Total pages: {num_pages}\n")
    
    # Create PowerPoint presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Process each page
    for page_idx in range(num_pages):
        print(f"{'='*60}")
        print(f"Processing page {page_idx + 1}/{num_pages}...")
        print(f"{'='*60}")
        
        # Step 1: Render PDF page as image
        print("  [1/3] Rendering PDF page...")
        page_image = render_pdf_page_as_image(pdf_path, page_idx, dpi)
        img_width, img_height = page_image.size
        
        # Step 2: Add slide with image
        print("  [2/3] Creating slide...")
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Calculate scaling to fit slide
        slide_aspect = prs.slide_width / prs.slide_height
        img_aspect = img_width / img_height
        
        if img_aspect > slide_aspect:
            pic_width = prs.slide_width
            pic_height = int(prs.slide_width * img_height / img_width)
            left = 0
            top = int((prs.slide_height - pic_height) / 2)
        else:
            pic_height = prs.slide_height
            pic_width = int(prs.slide_height * img_width / img_height)
            left = int((prs.slide_width - pic_width) / 2)
            top = 0
        
        # Add image to slide
        img_bytes = io.BytesIO()
        page_image.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        slide.shapes.add_picture(
            img_bytes,
            left,
            top,
            width=pic_width,
            height=pic_height
        )
        
        # Step 3: Generate and add speaker notes
        print("  [3/3] Generating speaker notes with AI...")
        notes = generate_speaker_notes(page_image, api_key)
        
        # Add notes to slide
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes
        
        print(f"  ✓ Slide {page_idx + 1} completed with notes")
        print(f"  Notes preview: {notes[:100]}...\n")
    
    # Save presentation
    print(f"{'='*60}")
    print(f"Saving presentation: {output_pptx}")
    prs.save(output_pptx)
    print("✓ Conversion completed successfully!")
    print("✓ Each slide has AI-generated speaker notes!")
    print(f"{'='*60}\n")
    
    return output_pptx


def main():
    """Main function."""
    
    if len(sys.argv) < 2:
        print("Usage: python pdf_to_pptx_with_notes.py <pdf_or_pptx_file> [output_pptx] [dpi]")
        print("\nExample:")
        print("  python pdf_to_pptx_with_notes.py presentation.pdf")
        print("  python pdf_to_pptx_with_notes.py presentation.pptx output.pptx")
        print("  python pdf_to_pptx_with_notes.py presentation.pdf output.pptx 200")
        print("\nSupports both PDF and PPTX input files.")
        print("Requires GOOGLE_API_KEY in .env file.")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    dpi = int(sys.argv[3]) if len(sys.argv) > 3 else 200
    
    # Check input file type
    file_ext = Path(input_file).suffix.lower()
    
    if not os.path.exists(input_file):
        print(f"✗ Error: File not found: {input_file}")
        sys.exit(1)
    
    try:
        if file_ext == '.pdf':
            # Process PDF
            result = pdf_to_pptx_with_notes(input_file, output_file, dpi)
        elif file_ext == '.pptx':
            # Process PPTX
            if output_file is None:
                output_file = str(Path(input_file).stem + '_with_notes.pptx')
            
            # Get API key
            api_key = os.environ.get('GOOGLE_API_KEY')
            if not api_key:
                raise ValueError("Google API key required. Set GOOGLE_API_KEY in .env file.")
            
            result = add_notes_to_pptx(input_file, output_file, api_key)
        else:
            print(f"✗ Error: Unsupported file format: {file_ext}")
            print("  Supported formats: .pdf, .pptx")
            sys.exit(1)
        
        print(f"\n✓ Successfully created: {result}")
    except Exception as e:
        print(f"\n✗ Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()


def process_pptx_with_progress(input_pptx, output_pptx, api_key, note_style="standard", note_tone="professional"):
    """
    Process PPTX with progress tracking for streaming updates.
    Yields progress dictionaries during processing.
    
    Args:
        input_pptx: Input PowerPoint file path
        output_pptx: Output PowerPoint file path
        api_key: Google AI API key
        note_style: Style of notes - 'brief', 'standard', or 'detailed'
        note_tone: Tone of notes - 'professional', 'casual', 'academic', or 'persuasive'
    """
    prs = Presentation(input_pptx)
    num_slides = len(prs.slides)
    
    yield {
        "status": "started",
        "total_slides": num_slides,
        "message": f"Starting to process {num_slides} slides..."
    }
    
    for idx, slide in enumerate(prs.slides):
        current_slide = idx + 1
        
        # Extract text content
        slide_text = []
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                slide_text.append(text)
        except Exception:
            pass
        
        combined_text = "\n".join(slide_text)
        
        # Try visual representation
        slide_image = render_slide_as_image(prs, idx)
        
        # Convert image to base64 for sending to UI
        slide_image_base64 = None
        if slide_image:
            try:
                # Resize image to reasonable size for web display (max 800px width)
                display_image = slide_image.copy()
                max_width = 800
                if display_image.width > max_width:
                    ratio = max_width / display_image.width
                    new_height = int(display_image.height * ratio)
                    display_image = display_image.resize((max_width, new_height), Image.LANCZOS)
                
                # Convert to base64
                img_buffer = io.BytesIO()
                display_image.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                slide_image_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
            except Exception:
                pass
        
        yield {
            "status": "processing",
            "current_slide": current_slide,
            "total_slides": num_slides,
            "message": f"Processing slide {current_slide} of {num_slides}...",
            "slide_image": slide_image_base64
        }
        
        notes = None
        
        if slide_image:
            try:
                notes = generate_speaker_notes(slide_image, api_key, note_style, note_tone)
            except Exception:
                notes = None
        
        # Fallback to text-only
        if notes is None and combined_text:
            try:
                notes = generate_notes_from_text(combined_text, api_key, note_style, note_tone)
            except Exception:
                notes = "This slide contains content but speaker notes could not be generated."
        
        if notes is None:
            notes = "This slide appears to be empty or contains only visual elements without text."
        
        # Add notes to slide
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
        except Exception:
            pass
        
        yield {
            "status": "processing",
            "current_slide": current_slide,
            "total_slides": num_slides,
            "message": f"Completed slide {current_slide} of {num_slides}",
            "slide_image": slide_image_base64
        }
    
    # Save presentation
    prs.save(output_pptx)
    
    yield {
        "status": "saving",
        "message": "Saving presentation..."
    }


def process_pdf_with_progress(pdf_path, output_pptx, dpi, api_key, note_style="standard", note_tone="professional"):
    """
    Process PDF with progress tracking for streaming updates.
    Yields progress dictionaries during processing.
    
    Args:
        pdf_path: Input PDF file path
        output_pptx: Output PowerPoint file path
        dpi: DPI for rendering
        api_key: Google AI API key
        note_style: Style of notes - 'brief', 'standard', or 'detailed'
        note_tone: Tone of notes - 'professional', 'casual', 'academic', or 'persuasive'
    """
    doc = fitz.open(pdf_path)
    num_pages = len(doc)
    doc.close()
    
    yield {
        "status": "started",
        "total_slides": num_pages,
        "message": f"Starting to convert {num_pages} pages..."
    }
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    for page_idx in range(num_pages):
        current_page = page_idx + 1
        
        # Render page
        page_image = render_pdf_page_as_image(pdf_path, page_idx, dpi)
        img_width, img_height = page_image.size
        
        # Convert image to base64 for sending to UI
        slide_image_base64 = None
        try:
            # Resize image to reasonable size for web display (max 800px width)
            display_image = page_image.copy()
            max_width = 800
            if display_image.width > max_width:
                ratio = max_width / display_image.width
                new_height = int(display_image.height * ratio)
                display_image = display_image.resize((max_width, new_height), Image.LANCZOS)
            
            # Convert to base64
            img_buffer = io.BytesIO()
            display_image.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            slide_image_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
        except Exception:
            pass
        
        yield {
            "status": "processing",
            "current_slide": current_page,
            "total_slides": num_pages,
            "message": f"Processing page {current_page} of {num_pages}...",
            "slide_image": slide_image_base64
        }
        
        # Add slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Calculate scaling
        slide_aspect = prs.slide_width / prs.slide_height
        img_aspect = img_width / img_height
        
        if img_aspect > slide_aspect:
            pic_width = prs.slide_width
            pic_height = int(prs.slide_width * img_height / img_width)
            left = 0
            top = int((prs.slide_height - pic_height) / 2)
        else:
            pic_height = prs.slide_height
            pic_width = int(prs.slide_height * img_width / img_height)
            left = int((prs.slide_width - pic_width) / 2)
            top = 0
        
        # Add image
        img_bytes = io.BytesIO()
        page_image.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        slide.shapes.add_picture(img_bytes, left, top, width=pic_width, height=pic_height)
        
        # Generate notes
        notes = generate_speaker_notes(page_image, api_key, note_style, note_tone)
        
        # Add notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes
        
        yield {
            "status": "processing",
            "current_slide": current_page,
            "total_slides": num_pages,
            "message": f"Completed page {current_page} of {num_pages}",
            "slide_image": slide_image_base64
        }
    
    # Save
    prs.save(output_pptx)
    
    yield {
        "status": "saving",
        "message": "Saving presentation..."
    }
